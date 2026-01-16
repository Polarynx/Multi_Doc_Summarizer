import os
import json
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

PROJECT_DOCS_DIR = "project_docs"
OUTPUT_FILE = "extracted_documents.json"


def load_docx(file_path):
    doc = Document(file_path)
    blocks = []

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:
            blocks.append({
                "source_file": os.path.basename(file_path),
                "source_type": "docx",
                "section": f"Paragraph {i + 1}",
                "content": text
            })

    return blocks


def load_pptx(file_path):
    presentation = Presentation(file_path)
    blocks = []

    for slide_idx, slide in enumerate(presentation.slides):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)

        if slide_text:
            blocks.append({
                "source_file": os.path.basename(file_path),
                "source_type": "pptx",
                "section": f"Slide {slide_idx + 1}",
                "content": "\n".join(slide_text)
            })

    return blocks


def load_xlsx(file_path):
    workbook = load_workbook(file_path, data_only=True)
    blocks = []

    for sheet in workbook.worksheets:
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            row_values = [str(cell) for cell in row if cell is not None]

            if row_values:
                blocks.append({
                    "source_file": os.path.basename(file_path),
                    "source_type": "xlsx",
                    "section": f"Sheet '{sheet.title}', Row {row_idx}",
                    "content": " | ".join(row_values)
                })

    return blocks


def extract_documents():
    all_blocks = []

    for filename in os.listdir(PROJECT_DOCS_DIR):
        file_path = os.path.join(PROJECT_DOCS_DIR, filename)

        if filename.lower().endswith(".docx"):
            all_blocks.extend(load_docx(file_path))

        elif filename.lower().endswith(".pptx"):
            all_blocks.extend(load_pptx(file_path))

        elif filename.lower().endswith(".xlsx"):
            all_blocks.extend(load_xlsx(file_path))

    return all_blocks


if __name__ == "__main__":
    extracted_data = extract_documents()

    with open(OUTPUT_FILE, "w", encoding="utf-8") as f:
        json.dump(extracted_data, f, indent=2, ensure_ascii=False)

    print(f"✅ Extracted {len(extracted_data)} text blocks")
    print(f"📄 Output written to {OUTPUT_FILE}")
